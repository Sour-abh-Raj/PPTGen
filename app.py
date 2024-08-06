import streamlit as st
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def adjust_font_size(shape, max_width, max_height, start_font_size=24):
    """
    Adjust the font size of the text in the shape to fit within the max_width and max_height.
    """
    current_font_size = start_font_size

    while True:
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.size = Pt(current_font_size)
        shape._element.getparent().append(shape._element)  # Re-render the shape
        if shape.width <= max_width and shape.height <= max_height:
            break
        current_font_size -= 1
        if current_font_size < 1:
            break
    return current_font_size

# Function to add slides to the presentation
def add_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    title_font_size = adjust_font_size(title_placeholder, slide_width, Inches(1.5), start_font_size=24)
    
    # Content
    content_placeholder = slide.placeholders[1]
    text_frame = content_placeholder.text_frame
    text_frame.clear()  # Clear any existing content
    
    content_font_size = title_font_size - 6 if title_font_size - 6 > 18 else 18
    for paragraph in content:
        p = text_frame.add_paragraph()
        p.text = paragraph
        p.font.size = Pt(content_font_size)
        p.alignment = PP_ALIGN.LEFT

    adjust_font_size(content_placeholder, slide_width, slide_height - Inches(1.5), start_font_size=content_font_size)

# Streamlit UI
st.title("Generate PowerPoint Presentation from JSON")

uploaded_json_file = st.file_uploader("Upload JSON file", type="json")
uploaded_template_file = st.file_uploader("Upload PowerPoint Template", type="pptx")

if uploaded_json_file is not None and uploaded_template_file is not None:
    # Read the JSON file
    slides_content = json.load(uploaded_json_file)

    # Load the template presentation
    prs = Presentation(uploaded_template_file)

    # Add slides to the presentation
    for slide_info in slides_content:
        add_slide(prs, slide_info["title"], slide_info["content"])

    # Optionally, delete the initial slides of the template
    while len(prs.slides) > len(slides_content):
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # Create output directory if it doesn't exist
    output_dir = "output_ppt"
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # Save the presentation
    output_file = os.path.join(output_dir, "Generated_Presentation.pptx")
    prs.save(output_file)

    st.success(f"Presentation generated successfully! Download it from the link below.")
    
    # Provide a download link for the generated PPT
    with open(output_file, "rb") as f:
        st.download_button(
            label="Download PPT",
            data=f,
            file_name="Generated_Presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

st.header("Made with ❤️ by Sourabh Raj")
